from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xD1, 0x00)
GRAY   = RGBColor(0xF2, 0xF4, 0xF8)
DARK   = RGBColor(0x2C, 0x2C, 0x2C)
ACCENT = RGBColor(0xE8, 0x4C, 0x3C)

W = Inches(13.33)
H = Inches(7.5)


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   size=20, color=DARK, marker="▶ "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = marker + item
        run.font.size = Pt(size)
        run.font.color.rgb = color


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, W, H, NAVY)
add_rect(s1, 0, Inches(5.2), W, Inches(2.3), RGBColor(0x14, 0x1E, 0x38))
add_rect(s1, Inches(0.5), Inches(1.6), Inches(0.08), Inches(3.2), YELLOW)

add_text(s1, "2026年10月、インボイスの「助け舟」が終わる",
         Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.4),
         size=36, bold=True, color=WHITE)
add_text(s1, "経理担当者が今すぐ動くべき理由",
         Inches(0.8), Inches(2.9), Inches(10), Inches(0.8),
         size=24, color=RGBColor(0xCC, 0xD6, 0xF0))
add_text(s1, "仕入税額控除が最大 50% 消える",
         Inches(0.8), Inches(4.0), Inches(10), Inches(0.7),
         size=20, color=YELLOW)
add_text(s1, "2026年  経理セミナー",
         Inches(0.8), Inches(5.5), Inches(6), Inches(0.6),
         size=16, color=RGBColor(0x99, 0xAA, 0xCC))

# ── Slide 2: Recap ──────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
add_rect(s2, 0, 0, W, Inches(1.3), NAVY)
add_rect(s2, 0, Inches(1.3), W, H - Inches(1.3), GRAY)

add_text(s2, "そもそもインボイス制度とは？  3分でおさらい",
         Inches(0.4), Inches(0.2), Inches(12), Inches(0.9),
         size=28, bold=True, color=WHITE)

cards = [
    ("適格請求書\n（インボイス）",   "登録番号・税率・税額が\n記載された請求書"),
    ("発行できるのは",               "税務署に登録した\n課税事業者のみ"),
    ("仕入税額控除",                 "インボイスなしでは\n消費税控除ができない"),
    ("免税事業者との取引",           "買い手側の\nコスト増につながる"),
]
cx = Inches(0.4)
for title, body in cards:
    add_rect(s2, cx, Inches(1.6), Inches(2.9), Inches(4.0), WHITE)
    add_rect(s2, cx, Inches(1.6), Inches(2.9), Inches(0.55), NAVY)
    add_text(s2, title, cx + Inches(0.12), Inches(1.65),
             Inches(2.65), Inches(0.5), size=14, bold=True, color=WHITE)
    add_text(s2, body,  cx + Inches(0.12), Inches(2.3),
             Inches(2.65), Inches(2.7), size=18, color=DARK)
    cx += Inches(3.1)

# ── Slide 3: Changes ────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
add_rect(s3, 0, 0, W, Inches(1.3), NAVY)
add_rect(s3, 0, Inches(1.3), W, H - Inches(1.3), GRAY)

add_text(s3, "2026年10月から何が変わるか",
         Inches(0.4), Inches(0.2), Inches(12), Inches(0.9),
         size=28, bold=True, color=WHITE)

rows = [
    ("〜 2026年9月",          "80%", NAVY),
    ("2026年10月〜2029年9月",  "50%", ACCENT),
    ("2029年10月〜",           " 0%", RGBColor(0x88,0x88,0x88)),
]
ry = Inches(1.55)
for period, pct, col in rows:
    add_rect(s3, Inches(0.4),  ry, Inches(6.0), Inches(1.3), WHITE)
    add_rect(s3, Inches(6.6),  ry, Inches(2.2), Inches(1.3), col)
    add_text(s3, period, Inches(0.55), ry + Inches(0.3),
             Inches(5.8), Inches(0.8), size=20, color=DARK)
    add_text(s3, pct + " 控除",
             Inches(6.65), ry + Inches(0.2),
             Inches(2.0), Inches(0.9), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ry += Inches(1.45)

add_rect(s3, Inches(9.2), Inches(1.55), Inches(3.7), Inches(4.2), WHITE)
add_text(s3, "具体例", Inches(9.35), Inches(1.7), Inches(3.4), Inches(0.5),
         size=14, bold=True, color=NAVY)
add_text(s3,
         "税込110万円の取引の場合\n\n"
         "現在（〜2026/9）\n控除： 8万円\n\n"
         "2026/10以降\n控除： 5万円\n\n"
         "→ 差額 3万円が\n   余分なコストに",
         Inches(9.35), Inches(2.2), Inches(3.4), Inches(3.3),
         size=15, color=DARK)

# ── Slide 4: Checklist ──────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
add_rect(s4, 0, 0, W, Inches(1.3), NAVY)
add_rect(s4, 0, Inches(1.3), W, H - Inches(1.3), GRAY)

add_text(s4, "今すぐ確認！自社の対応状況チェックリスト",
         Inches(0.4), Inches(0.2), Inches(12), Inches(0.9),
         size=28, bold=True, color=WHITE)

checks_l = [
    "仕入先・外注先の登録番号リストを整備している",
    "未登録の取引先を把握し、コスト試算を済ませた",
    "2026年10月以降の税負担増を計算し経営判断に反映した",
]
checks_r = [
    "受領した適格請求書を適切に保存・管理している",
    "会計ソフトがインボイス対応の設定になっている",
    "発行請求書に登録番号・税率・税額が記載されている",
]

def add_checklist(slide, items, left, top):
    for i, item in enumerate(items):
        y = top + Inches(i * 1.55)
        add_rect(slide, left, y, Inches(5.8), Inches(1.35), WHITE)
        add_rect(slide, left, y, Inches(0.5), Inches(1.35), NAVY)
        add_text(slide, "□", left + Inches(0.05), y + Inches(0.3),
                 Inches(0.4), Inches(0.7), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, item, left + Inches(0.6), y + Inches(0.2),
                 Inches(5.1), Inches(1.0), size=17, color=DARK)

add_checklist(s4, checks_l, Inches(0.4),  Inches(1.55))
add_checklist(s4, checks_r, Inches(6.85), Inches(1.55))

# ── Slide 5: Action Steps ───────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank)
add_rect(s5, 0, 0, W, Inches(1.3), NAVY)
add_rect(s5, 0, Inches(1.3), W, H - Inches(1.3), GRAY)

add_text(s5, "今日から始める 3ステップ対応",
         Inches(0.4), Inches(0.2), Inches(12), Inches(0.9),
         size=28, bold=True, color=WHITE)

steps = [
    ("Step 1", "今週中",
     "取引先リストの棚卸し\n国税庁サイトで登録番号を確認\n未登録先をリストアップ"),
    ("Step 2", "今月中",
     "コスト試算と対応方針の決定\n50%控除時の税負担増を計算\n取引継続 / 価格交渉 / 代替先を検討"),
    ("Step 3", "2026年9月末まで",
     "システム・フロー整備\n会計ソフト設定の更新\n請求書テンプレートに登録番号を追加"),
]
sx = Inches(0.4)
for label, timing, body in steps:
    add_rect(s5, sx, Inches(1.55), Inches(3.9), Inches(5.4), WHITE)
    add_rect(s5, sx, Inches(1.55), Inches(3.9), Inches(0.8), NAVY)
    add_text(s5, label,  sx + Inches(0.15), Inches(1.6),
             Inches(1.5), Inches(0.65), size=18, bold=True, color=WHITE)
    add_text(s5, timing, sx + Inches(1.7),  Inches(1.65),
             Inches(2.0), Inches(0.55), size=14, color=YELLOW)
    add_text(s5, body,   sx + Inches(0.2),  Inches(2.55),
             Inches(3.5), Inches(4.0), size=17, color=DARK)
    sx += Inches(4.3)

out = "/home/user/anthropic-academy/invoice_seminar_slides.pptx"
prs.save(out)
print(f"Saved: {out}")
